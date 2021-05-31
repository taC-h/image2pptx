#python3.8.3
import itertools
import os
import pathlib
import shutil

#IO
from pptx import Presentation
from PIL import Image
from pdf2image import convert_from_path
import json
from io import BytesIO
import subprocess
from glob import glob

#GUI
import tkinter as tk
import tkinter.filedialog as dialog
from tkinter.messagebox import showinfo,askretrycancel


class Application(tk.Frame):
    input_ftype = [
        ("全てのファイル","*"),
        ("JPEG",[
            "*.jpeg",
            "*.jpg",
            "*.jpe",
            "*.jfif",
            ],
        ),
        ("ビットマップ","*bmp"),
        ("GIF","*.gif",),
        ("TIFF",[
            "*.tif",
            "*.tiff",
            ],
        ),
        ("PNG","*.png",),
        ("HEIC","*.heic",),
        ("PDF","*.pdf"),
        ]
    iDir = os.path.dirname(__file__)
    pptx_ftype = [("presentation","*.pptx")]
    pdf_ftype = [("PDF","*.pdf")]
    heic_ftype = [("HEIC","*.heic")]
    config_template = {
        "dpi":150,
        "path":{
            "pdfdir":"",
            "imagedir":"",
            "slidedir":"",
        },
        "image_save_ext":".png",
        "pathrewrite":True,
        }
    name = "image2pptx"

    poppler_dir = os.path.join(pathlib.Path(iDir), "source\\poppler\\bin")
    magick_dir = os.path.join(pathlib.Path(iDir), "source\\ImageMagick")
    heic_dir = os.path.join(pathlib.Path(iDir), "tmp")
    os.environ["PATH"] += os.pathsep + poppler_dir
    if not glob(heic_dir):
        os.mkdir(heic_dir)


    def __init__(self, master=None):
        super().__init__(master)
        self.template_path = os.path.join(self.iDir, "source/template.pptx")
        self.configfile = os.path.join(self.iDir,"config.json")


        self.Get_config()
        self.master = master
        self.pack()
        self.create_widgets()


    def create_widgets(self):
        self.image_button = tk.Button(self, text="image2pptx", command=self.image2pptx)
        self.image_button.pack()

        self.pdf_button = tk.Button(self, text="pdf2image", command=self.pdf2image)
        self.pdf_button.pack()

        self.heic_button = tk.Button(self, text="heic2image",command=self.heic2image)
        self.heic_button.pack()

        self.qu = tk.Button(self, text="QUIT", fg="red", command=self.__quit)
        self.qu.pack(side="bottom")


    def image2pptx(self):
        self.prs = Presentation(self.template_path)
        self.blank_slide_layout = self.prs.slide_layouts[6]
        self.aspect:float = self.prs.slide_width / self.prs.slide_height

        if not os.path.isdir(self.config["path"]["imagedir"]):
            self.config["path"]["imagedir"] = self.iDir
        file_list = list(dialog.askopenfilenames(filetypes = self.input_ftype,initialdir = self.config["path"]["imagedir"],title="ファイルを選択してください"))
        if not file_list:
            showinfo(self.name,"画像が選択されませんでした")
            return
        self.pathrewrite("imagedir",os.path.dirname(file_list[0]))

        try:
            self.files_classifying2pptx(file_list)
        except :
            showinfo("image2pptx","変換に失敗しました")

        self.tmpremove()

        while True:
            if not os.path.isdir(self.config["path"]["slidedir"]):
                self.config["path"]["slidedir"] = self.iDir
            ret = dialog.asksaveasfilename(defaultextension="pptx",filetypes= self.pptx_ftype,initialdir = self.config["path"]["slidedir"],title="保存するファイル名を決めてください")
            if ret:
                break
            else:
                res = askretrycancel("image2pptx","ファイル名が入力されませんでした")
            if not res:
                return

        try:
            self.prs.save(ret)
        except:
            showinfo("image2pptx","保存に失敗しました")
        else:
            self.pathrewrite("slidedir",os.path.dirname(ret))
        #fin


    def pdf2image(self):
        if not os.path.isdir(self.config["path"]["pdfdir"]):
            self.config["path"]["pdfdir"] = self.iDir
        file_list = list(dialog.askopenfilenames(filetypes = self.pdf_ftype,initialdir = self.config["path"]["pdfdir"],title="ファイルを選択してください"))
        if not file_list:
            showinfo(self.name,"ファイルが選択されませんでした")
            return
        self.pathrewrite("pdfdir",os.path.dirname(file_list[0]))

        filename_list = [os.path.splitext(os.path.basename(f))[1] for f in file_list]
        try:
            images = [self.PDF_Converter(f) for f in file_list]
        except:
            showinfo("pdf2image","変換に失敗しました")
            return

        if not os.path.isdir(self.config["path"]["imagedir"]):
            self.config["path"]["imagedir"] = self.iDir


        while True:
            if not os.path.isdir(self.config["path"]["slidedir"]):
                self.config["path"]["slidedir"] = self.iDir
            save_dir:str = dialog.askdirectory(initialdir = self.config["path"]["imagedir"],title="保存するフォルダを選択してください")
            if save_dir:
                break
            else:
                res = askretrycancel("image2pptx","フォルダが選択されませんでした")
            if not res:
                return

        try:
            for i, filename in zip(images, filename_list):
                for pages, num in zip(i, itertools.count(1)):
                    save_path = os.path.join(save_dir,"{}_{:02d}{}".format(filename, num,self.config["image_save_ext"]))
                    pages.save(save_path)
        except:
            showinfo("pdf2image","保存に失敗しました")
        else:
            self.pathrewrite("imagedir",os.path.dirname(save_dir))


    def heic2image(self):
        if not os.path.isdir(self.config["path"]["imagedir"]):
            self.config["path"]["imagedir"] = self.iDir
        file_list:list[str] = list(dialog.askopenfilenames(filetypes = self.heic_ftype,initialdir = self.config["path"]["imagedir"],title="ファイルを選択してください"))
        if not file_list:
            showinfo("heic2image","ファイルが選択されませんでした")
            return

        self.pathrewrite("imagedir",os.path.dirname(file_list[0]))

        ret_list=[]
        for f in file_list:
            ret_list.append(self.HEIC_Converter(f))

        while True:
            if not os.path.isdir(self.config["path"]["slidedir"]):
                self.config["path"]["slidedir"] = self.iDir
            save_dir:str = dialog.askdirectory(initialdir = self.config["path"]["imagedir"],title="保存するフォルダを選択してください")
            if save_dir:
                break
            else:
                res = askretrycancel("heic2image","フォルダが選択されませんでした")
            if not res:
                return

        #save
        try:
            for f in ret_list:
                save_name = os.path.join(save_dir, os.path.basename(f))
                shutil.copy(f,save_dir)
        except:
            showinfo("heic2image","保存に失敗しました")
        self.tmpremove()

    def __quit(self):
        try:
            with open(self.configfile,"w") as f:
                json.dump(self.config,f,indent=4,ensure_ascii=False)
        except:
            showinfo(self.name,"configファイルが更新されませんでした")

        self.tmpremove()
        self.master.quit()
        self.master.destroy()


    def Get_config(self):
        try:
            with open(self.configfile) as f:
                self.config:dict = json.load(f)
        except:
            self.config = self.config_template.copy()
            self.config["path"]["imagedir"] = self.config["path"]["pdfdir"] = self.config["path"]["slidedir"] = self.iDir


    def files_classifying2pptx(self,file_list:list):
        addedpict = ""#has attribute read
        for f in file_list:
            ftype = type(f)
            if ftype in [list,tuple]:
                self.files_classifying2pptx(f)
                continue
            elif ftype == str:
                ext = os.path.splitext(f)[1]

                if ext == ".pdf":
                    self.files_classifying2pptx(self.PDF_Converter(f))
                    continue
                elif ext == ".heic":
                    addedpict = self.HEIC_Converter(f)
                    im = Image.open(addedpict)
                else:
                    addedpict = f
                    im = Image.open(f)
            else:#from pdf
                tmp = BytesIO()
                f.save(tmp,"JPEG")
                addedpict = tmp
                im =  Image.open(tmp)

            im_aspect = im.size[0] / im.size[1]
            slide = self.prs.slides.add_slide(self.blank_slide_layout)
            if(im_aspect > self.aspect):
                pic_height = self.prs.slide_width /im_aspect
                pic = slide.shapes.add_picture(addedpict, 0, (self.prs.slide_height - pic_height)/2, self.prs.slide_width, pic_height)
            else:
                pic_width = self.prs.slide_height * im_aspect
                pic = slide.shapes.add_picture(addedpict, (self.prs.slide_width - pic_width)/2 , 0, pic_width, self.prs.slide_height)
            im.close()


    def PDF_Converter(self, filename:str) -> list:
        """return [imageobject,]"""
        return convert_from_path(filename,self.config["dpi"])


    def HEIC_Converter(self, filename:str) -> str:
        """return filepath"""
        heicname = os.path.basename(filename)
        tmpheic = os.path.join(self.heic_dir,heicname)#fullpath
        shutil.copy(filename,tmpheic)
        ret = os.path.splitext(tmpheic)[0] + self.config["image_save_ext"]#fullpath
        converer = os.path.join(self.magick_dir,"converting.bat")
        with open(converer,mode="w") as f:
            f.write("{} {} {}".format(os.path.join(self.magick_dir,"convert"), tmpheic, ret,))
        subprocess.call(f"{converer}",shell=True)
        return ret#fullpath


    def pathrewrite(self, string:str, dirname:str):
        if self.config["pathrewrite"]:
            self.config["path"][string] = dirname


    def tmpremove(self) -> None:
        try:
            shutil.rmtree(self.heic_dir)
            os.mkdir(self.heic_dir)
        except:
            pass



if __name__ == "__main__":
    root = tk.Tk()
    root.title("image2pptx2.0")
    root.geometry("200x150")
    app = Application(master=root)
    app.mainloop()